"""Extract text + structured data from PDF, Excel, Word, PowerPoint."""

import io
from dataclasses import dataclass, field


@dataclass
class ExtractedDocument:
    name: str
    file_type: str
    text_chunks: list[str] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)  # list of tables as 2D arrays
    metadata: dict = field(default_factory=dict)


def extract(filename: str, raw_bytes: bytes) -> ExtractedDocument:
    ext = filename.rsplit(".", 1)[-1].lower()
    doc = ExtractedDocument(name=filename, file_type=ext)

    if ext == "pdf":
        _extract_pdf(raw_bytes, doc)
    elif ext in ("xlsx", "xls", "csv"):
        _extract_excel(raw_bytes, doc, ext)
    elif ext in ("docx", "doc"):
        _extract_word(raw_bytes, doc)
    elif ext in ("pptx", "ppt"):
        _extract_pptx(raw_bytes, doc)
    else:
        doc.text_chunks = [raw_bytes.decode("utf-8", errors="replace")]

    return doc


def _extract_pdf(raw: bytes, doc: ExtractedDocument) -> None:
    import pdfplumber

    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                doc.text_chunks.append(text)
            tables = page.extract_tables() or []
            for table in tables:
                clean = [[cell or "" for cell in row] for row in table]
                doc.tables.append(clean)
    doc.metadata["pages"] = len(pdfplumber.open(io.BytesIO(raw)).pages)


def _extract_excel(raw: bytes, doc: ExtractedDocument, ext: str) -> None:
    import pandas as pd

    if ext == "csv":
        df = pd.read_csv(io.BytesIO(raw))
        doc.tables.append([list(df.columns)] + df.astype(str).values.tolist())
        doc.text_chunks.append(df.to_string(index=False))
    else:
        xl = pd.ExcelFile(io.BytesIO(raw))
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            doc.tables.append([list(df.columns)] + df.astype(str).values.tolist())
            doc.text_chunks.append(f"Sheet: {sheet}\n{df.to_string(index=False)}")
    doc.metadata["sheets"] = len(doc.tables)


def _extract_word(raw: bytes, doc: ExtractedDocument) -> None:
    from docx import Document

    d = Document(io.BytesIO(raw))
    paragraphs = [p.text for p in d.paragraphs if p.text.strip()]
    doc.text_chunks = paragraphs
    for table in d.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        doc.tables.append(rows)


def _extract_pptx(raw: bytes, doc: ExtractedDocument) -> None:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        if texts:
            doc.text_chunks.append(f"Slide {slide_num}:\n" + "\n".join(texts))
